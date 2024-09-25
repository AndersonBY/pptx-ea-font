import unittest
from typing import Optional, Dict

import pptx_ea_font
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx import Presentation


def find_with_optional_ns(element, path: str, namespaces: Optional[Dict[str, str]] = None):
    return element.find(path, namespaces=namespaces)


class TestSetFont(unittest.TestCase):
    def test_set_font(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(7.5))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = "测试文本"

        pptx_ea_font.set_font(run, "微软雅黑")

        self.assertEqual(run.font.name, "微软雅黑")
        ea_element = find_with_optional_ns(run.font._rPr, qn("a:ea"))
        self.assertIsNotNone(ea_element)
        self.assertEqual(ea_element.get("typeface"), "微软雅黑")

    def test_change_font_to_kaiti(self):
        prs = Presentation("tests/test.pptx")
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_ea_font.set_font(run, "楷体")

        prs.save("tests/test_kaiti.pptx")


if __name__ == "__main__":
    unittest.main()
